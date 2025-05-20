from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define a function to add a slide with a title and content
def add_slide(title, content, highlights=None, colors=None):
    slide_layout = prs.slide_layouts[1]  # Use the "Title and Content" layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    # Set the title
    title_placeholder.text = title

    # Set the content with optional highlights
    text_frame = content_placeholder.text_frame
    p = text_frame.add_paragraph()
    p.text = content
    p.font.size = Pt(24)
    
    # Apply highlights if any
    if highlights:
        for highlight, color in zip(highlights, colors):
            if highlight in p.text:
                start_idx = p.text.index(highlight)
                end_idx = start_idx + len(highlight)
                p.text = p.text.replace(highlight, "")
                run = p.add_run()
                run.text = highlight
                run.font.bold = True
                run.font.color.rgb = color

# Slide 1: Obiettivo
add_slide("Obiettivo", "")

# Slide 2: Dato
content_dato = (
    "max 10x₁ + 7x₂ + 9x₃\n"
    "2x₁ + 3x₂ + 2x₃ = 5\n"
    "3x₁ + 2x₂ + 3x₃ ≤ 5\n"
    "2x₁ + 3x₂ + x₃ ≥ 3"
)
highlights_dato = ["3x₁ + 2x₂ + 3x₃ ≤ 5", "2x₁ + 3x₂ + x₃ ≥ 3"]
colors_dato = [RGBColor(0, 176, 240), RGBColor(0, 176, 80)]
add_slide("Dato", content_dato, highlights_dato, colors_dato)

# Slide 3: Trasformarlo in un problema QUBO
add_slide("Trasformarlo in un problema QUBO", "")

# Slide 4: Difficoltà
content_difficulty = (
    "Trasformare (1) (2) in vincoli equazionali\n"
    "minimizzando il numero di variabili per passare a QUBO"
)
highlights_difficulty = ["(1)", "(2)", "minimizzando"]
colors_difficulty = [RGBColor(255, 0, 0), RGBColor(255, 0, 0), RGBColor(255, 0, 0)]
add_slide("Difficoltà", content_difficulty, highlights_difficulty, colors_difficulty)

# Save the presentation
pptx_file_corrected = ".\Transformation_QUBO_Presentation_Corrected.pptx"
prs.save(pptx_file_corrected)

# pptx_file_corrected
